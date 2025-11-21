from pptx import Presentation
from pptx.util import Inches, Pt
import os

IMAGES = {
    "Distribución morosidad": "PowerBI/plot_morosidad.png",
    "AUC comparativo": "PowerBI/plot_auc_comparativa.png",
    "Curva ROC comparativa": "PowerBI/roc_comparativa.png",
    "Matriz confusión mejor modelo": "PowerBI/matriz_confusion_mejor_modelo.png",
    "Importancia de variables": "PowerBI/importancia_variables.xlsx"
}

prs = Presentation()

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_text_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def add_image_slide(title, image_path):
    if not os.path.exists(image_path):
        return
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))

add_title_slide("Informe Analítico - Morosidad Crediticia",
                "Generado automáticamente desde Python")

add_text_slide("Objetivo del Proyecto",
"""
Analizar factores socioeconómicos que explican la morosidad y construir modelos predictivos
para estimar el riesgo de incumplimiento de pago.
""")

add_image_slide("Distribución de Morosidad", "PowerBI/plot_morosidad.png")
add_image_slide("Comparación AUC de Modelos", "PowerBI/plot_auc_comparativa.png")
add_image_slide("Curvas ROC Comparativas", "PowerBI/roc_comparativa.png")
add_image_slide("Matriz de Confusión del Mejor Modelo", "PowerBI/matriz_confusion_mejor_modelo.png")

add_text_slide("Conclusiones",
"""
• El modelo con mejor desempeño fue Random Forest.
• La variable 'relación deuda/ingresos' es altamente predictiva.
• Los clientes de mayor riesgo se concentran en ingresos bajos y alta antigüedad crediticia.
• La clasificación de riesgo (bajo/medio/alto) permite segmentación inmediata del portafolio.
""")

add_text_slide("Recomendaciones",
"""
• Implementar Random Forest en el proceso de scoring.
• Monitorear clientes con score > 66%.
• Incorporar más variables de comportamiento para mejorar predicción.
• Desarrollar alertas preventivas basadas en el score de riesgo.
""")

prs.save("PowerBI/Informe_Morosidad.pptx")
print("PowerPoint generado exitosamente en PowerBI/Informe_Morosidad.pptx")
